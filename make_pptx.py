#!/usr/bin/env python3
"""Generate 3 PPTX guides: client, nanny, admin"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

SC = "/home/user/nanny_app/screenshots/resized"
OUT = "/home/user/nanny_app"

TEAL   = RGBColor(0x01, 0x6b, 0x82)
GOLD   = RGBColor(0xD4, 0xAF, 0x37)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1a, 0x1a, 0x2e)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
GREEN  = RGBColor(0x2e, 0xa0, 0x43)

W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, text, size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        bg_color=None, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    if bg_color:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg_color
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def img(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    from PIL import Image as PILImage
    im = PILImage.open(path)
    iw, ih = im.size
    if h is None:
        h = int(w * ih / iw)
    slide.shapes.add_picture(path, l, t, w, h)

def rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def cover_slide(prs, title, subtitle, accent, sc_path=None):
    sl = blank(prs)
    bg(sl, DARK)
    rect(sl, 0, 0, W, Inches(0.1), accent)
    rect(sl, 0, H - Inches(0.08), W, Inches(0.08), accent)
    if sc_path and os.path.exists(sc_path):
        img(sl, sc_path, Inches(7.5), Inches(0.5), Inches(5.5))
        rect(sl, Inches(7.3), 0, Inches(0.2), H, RGBColor(0x1a,0x1a,0x2e))
    box(sl, Inches(0.6), Inches(1.5), Inches(6.5), Inches(1.2), title,
        size=40, bold=True, color=WHITE)
    box(sl, Inches(0.6), Inches(2.9), Inches(6.5), Inches(0.6), subtitle,
        size=22, color=accent)
    box(sl, Inches(0.6), Inches(6.5), Inches(6), Inches(0.5),
        "Nanny Nha Trang • nannynhatrang.ru", size=14, color=RGBColor(0x88,0x88,0x88))

def section_slide(prs, num, title, accent):
    sl = blank(prs)
    bg(sl, accent)
    box(sl, Inches(0.8), Inches(2.5), Inches(11), Inches(1.2),
        f"{num}. {title}", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def content_slide(prs, title, bullets, sc_path=None, accent=TEAL):
    sl = blank(prs)
    bg(sl, LGRAY)
    rect(sl, 0, 0, W, Inches(1.1), accent)
    box(sl, Inches(0.4), Inches(0.15), Inches(12), Inches(0.8),
        title, size=26, bold=True, color=WHITE)
    text_w = Inches(6.5) if sc_path else Inches(12.5)
    y = Inches(1.4)
    for b in bullets:
        if b.startswith("##"):
            box(sl, Inches(0.5), y, text_w, Inches(0.45),
                b[2:].strip(), size=15, bold=True, color=accent)
            y += Inches(0.5)
        else:
            box(sl, Inches(0.7), y, text_w, Inches(0.5),
                b, size=14, color=DARK)
            y += Inches(0.45)
    if sc_path and os.path.exists(sc_path):
        img(sl, sc_path, Inches(7.3), Inches(1.2), Inches(5.7))

def step_slide(prs, step_num, title, desc, sc_path=None, accent=TEAL):
    sl = blank(prs)
    bg(sl, WHITE)
    rect(sl, 0, 0, Inches(1.2), H, accent)
    box(sl, Inches(0.1), Inches(2.8), Inches(1.0), Inches(1.5),
        str(step_num), size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sl, Inches(1.4), Inches(0.3), Inches(11), Inches(0.8),
        title, size=30, bold=True, color=accent)
    box(sl, Inches(1.4), Inches(1.3), Inches(5.5), Inches(5.5),
        desc, size=16, color=DARK, wrap=True)
    if sc_path and os.path.exists(sc_path):
        img(sl, sc_path, Inches(7.3), Inches(1.2), Inches(5.7))

def tip_slide(prs, title, tips, color=GREEN):
    sl = blank(prs)
    bg(sl, color)
    box(sl, Inches(0.6), Inches(0.4), Inches(12), Inches(0.9),
        "💡 " + title, size=30, bold=True, color=WHITE)
    y = Inches(1.5)
    for t in tips:
        rect(sl, Inches(0.5), y, Inches(12.3), Inches(0.75), RGBColor(0xff,0xff,0xff))
        box(sl, Inches(0.7), y+Inches(0.1), Inches(12), Inches(0.55),
            "✓ " + t, size=15, color=color, bold=True)
        y += Inches(0.88)


# ═══════════════════════════════════════
# ПРЕЗЕНТАЦИЯ 1 — КЛИЕНТ
# ═══════════════════════════════════════
def make_client():
    prs = new_prs()
    cover_slide(prs, "Личный кабинет\nКлиента",
                "Как пользоваться порталом Nanny Nha Trang", GOLD,
                f"{SC}/sc_client1.png")

    section_slide(prs, 1, "Начало работы", TEAL)

    step_slide(prs, 1, "Откройте Mini App в Telegram",
        "1. Найдите бота @Nannynhatrang_bot\n"
        "2. Нажмите кнопку Menu или Start\n"
        "3. Откроется форма заявки\n\n"
        "Telegram автоматически вас идентифицирует — \n"
        "никаких паролей не нужно.",
        f"{SC}/sc_tg_entry.png")

    step_slide(prs, 2, "Заполните заявку",
        "Укажите:\n"
        "• Ваше имя\n"
        "• Telegram для связи\n"
        "• Имя и возраст ребёнка\n"
        "• Желаемые даты\n"
        "• Время начала и продолжительность\n\n"
        "Нажмите «Отправить заявку».\n"
        "Вы получите ссылку на личный кабинет.",
        f"{SC}/sc_home2.png")

    step_slide(prs, 3, "Ваш личный кабинет",
        "После отправки заявки вы получите:\n"
        "• Личную ссылку /client/ваш-токен\n"
        "• Статус заявки\n"
        "• Информацию о назначенной няне\n"
        "• Календарь смен\n"
        "• Документы и договоры\n\n"
        "Сохраните ссылку — она ваша навсегда.",
        f"{SC}/sc_client1.png")

    section_slide(prs, 2, "Функции кабинета", TEAL)

    content_slide(prs, "Календарь смен", [
        "## Как читать календарь",
        "🟢 Зелёные даты — подтверждённые смены",
        "🟡 Жёлтые даты — встреча с няней",
        "⬜ Обычные даты — доступны для бронирования",
        "",
        "## Дополнительно",
        "Нажмите на дату смены — откроются чеки",
        "Смены синхронизированы с порталом няни",
        "Обновляются в реальном времени",
    ], f"{SC}/sc_client2.png")

    content_slide(prs, "Документы и договоры", [
        "## Доступные документы",
        "📄 Договор агентство–родители",
        "📄 Договор няня–родители",
        "📦 Пакет документов для клиентов",
        "",
        "## Тарифы и программы",
        "💰 Инфографика с тарифами агентства",
        "🤝 Реферальная программа",
        "",
        "Все документы открываются через Telegram",
    ], f"{SC}/sc_client1.png")

    content_slide(prs, "Профиль няни", [
        "## Кнопка «Профиль няни →»",
        "Открывает страницу вашей закреплённой няни",
        "Фото, опыт работы, языки",
        "Рейтинг и отзывы",
        "",
        "## Статус заявки",
        "✅ Назначена — чётко видно в шапке ЛК",
        "⏳ Подбираем — обычно несколько часов",
        "Вы получите уведомление в Telegram",
    ], f"{SC}/sc_client1.png")

    tip_slide(prs, "Советы клиенту", [
        "Сохраните ссылку на ЛК в закладки Telegram",
        "Следите за уведомлениями — они приходят в бот",
        "При вопросах пишите менеджеру в @Nannynhatrang_bot",
        "Фото-чеки можно загрузить прямо в ЛК",
        "Статьи в блоге помогут выбрать правильный тариф",
    ], TEAL)

    prs.save(f"{OUT}/guide_client.pptx")
    print("✅ guide_client.pptx")


# ═══════════════════════════════════════
# ПРЕЗЕНТАЦИЯ 2 — НЯНЯ
# ═══════════════════════════════════════
def make_nanny():
    prs = new_prs()
    cover_slide(prs, "Портал Няни",
                "Инструкция по работе с порталом Nanny Nha Trang", GREEN,
                f"{SC}/sc_nanny1.png")

    section_slide(prs, 1, "Ваш личный портал", GREEN)

    step_slide(prs, 1, "Ссылка на ваш портал",
        "Администратор выдаёт вам персональную ссылку:\n"
        "/nanny/ваш-токен\n\n"
        "Или через Telegram бот @Nannynhatrang_bot\n"
        "после добавления вас в систему.\n\n"
        "Ссылка работает без пароля — \n"
        "храните её как конфиденциальную.",
        f"{SC}/sc_nanny1.png", GREEN)

    step_slide(prs, 2, "Что вы видите в портале",
        "• Ваше фото и профиль\n"
        "• Список назначенных смен\n"
        "• Календарь занятости\n"
        "• Блоки выходных дней\n"
        "• Возможность загрузить чек после смены\n\n"
        "Портал обновляется при каждом назначении.",
        f"{SC}/sc_nanny1.png", GREEN)

    section_slide(prs, 2, "Управление сменами", GREEN)

    content_slide(prs, "Ваш календарь", [
        "## Цвета в календаре",
        "🟢 Рабочий день — назначена смена",
        "🔴 Выходной — вы отметили как недоступный",
        "⬜ Свободный день",
        "",
        "## Как отметить выходной",
        "Нажмите на дату в календаре",
        "Выберите «Отметить выходной»",
        "Дата станет недоступной для назначений",
    ], f"{SC}/sc_nanny1.png", GREEN)

    content_slide(prs, "Загрузка чеков", [
        "## После каждой смены",
        "1. Сфотографируйте подтверждение оплаты",
        "2. Зайдите на свой портал /nanny/токен",
        "3. Найдите нужную дату в списке смен",
        "4. Нажмите «Загрузить чек»",
        "5. Выберите фото",
        "",
        "Чек будет виден и клиенту, и администратору",
        "Уведомление придёт автоматически",
    ], f"{SC}/sc_nanny1.png", GREEN)

    section_slide(prs, 3, "Уведомления", GREEN)

    content_slide(prs, "Когда приходят уведомления в Telegram", [
        "## Вы получаете уведомление когда:",
        "🆕 Администратор назначил вас на новую заявку",
        "📅 Клиент подтвердил даты смен",
        "💰 Клиент загрузил чек об оплате",
        "⏰ За 2 часа до начала смены (напоминание)",
        "",
        "## Для этого нужно:",
        "Ваш Telegram ID прописан в профиле в системе",
        "Обратитесь к администратору если нет уведомлений",
    ], f"{SC}/sc_nanny1.png", GREEN)

    tip_slide(prs, "Советы для няни", [
        "Заходите в портал накануне — проверяйте расписание",
        "Отмечайте выходные заранее чтобы не было конфликтов",
        "Загружайте чеки в тот же день после смены",
        "При вопросах пишите администратору в Telegram",
        "Ваш профиль виден клиентам — держите его актуальным",
    ], GREEN)

    prs.save(f"{OUT}/guide_nanny.pptx")
    print("✅ guide_nanny.pptx")


# ═══════════════════════════════════════
# ПРЕЗЕНТАЦИЯ 3 — АДМИН
# ═══════════════════════════════════════
def make_admin():
    prs = new_prs()
    cover_slide(prs, "Административная\nПанель",
                "Полное руководство по управлению агентством", RGBColor(0x7c,0x3a,0xed),
                f"{SC}/sc_admin_login.png")

    section_slide(prs, 1, "Вход в систему", RGBColor(0x7c,0x3a,0xed))

    step_slide(prs, 1, "Авторизация через Telegram",
        "1. Откройте Mini App @Nannynhatrang_bot\n"
        "2. Система определяет ваш Telegram ID\n"
        "3. Если ID в списке ADMIN_IDS → вы в админке\n\n"
        "Через браузер (разработка):\n"
        "/admin?dev_token=ВАШ_ADMIN_TOKEN\n\n"
        "Railway: переменные ADMIN_IDS, ADMIN_TOKEN,\n"
        "TELEGRAM_BOT_TOKEN, FLASK_SECRET_KEY",
        f"{SC}/sc_admin_login.png", RGBColor(0x7c,0x3a,0xed))

    section_slide(prs, 2, "Управление заявками", RGBColor(0x7c,0x3a,0xed))

    content_slide(prs, "Список заявок", [
        "## Вкладка «Заявки»",
        "Все входящие заявки от клиентов",
        "Поиск по имени, телефону, имени ребёнка",
        "Статус: ожидает назначения / назначена",
        "",
        "## Действия",
        "Назначить няню — выпадающий список",
        "Открыть ЛК клиента — кнопка «ЛК →»",
        "Все данные в одном месте",
    ], f"{SC}/sc_desktop_home.png", RGBColor(0x7c,0x3a,0xed))

    step_slide(prs, 2, "Назначение няни на заявку",
        "1. Найдите заявку в разделе «Заявки»\n"
        "2. В колонке «Няня» выберите из списка\n"
        "3. Нажмите «Сохранить»\n\n"
        "Автоматически:\n"
        "• Няня получит уведомление в Telegram\n"
        "• Клиент получит уведомление с ссылкой на ЛК\n"
        "• Календарь обновится у обоих\n\n"
        "Система проверяет конфликты расписания!",
        f"{SC}/sc_desktop_home.png", RGBColor(0x7c,0x3a,0xed))

    section_slide(prs, 3, "Управление нянями", RGBColor(0x7c,0x3a,0xed))

    content_slide(prs, "Добавление и редактирование няни", [
        "## Раздел «Няни»",
        "Список всех зарегистрированных нянь",
        "",
        "## При добавлении/редактировании:",
        "• Имя, фото, возраст, опыт",
        "• Telegram ID (для уведомлений)",
        "• Портальный токен (уникальная ссылка)",
        "• Языки, специализация",
        "",
        "## Важно:",
        "Telegram ID — числовой, не @username",
        "Portal token — латиница, без пробелов",
    ], f"{SC}/sc_nanny1.png", RGBColor(0x7c,0x3a,0xed))

    section_slide(prs, 4, "Статьи и блог", RGBColor(0x7c,0x3a,0xed))

    step_slide(prs, 4, "Создание статьи",
        "1. Раздел «Статьи» → «+ Новая статья»\n"
        "2. Заполните:\n"
        "   • Заголовок (обязательно)\n"
        "   • Slug — URL статьи /blog/slug\n"
        "   • Краткое описание\n"
        "   • Тело статьи в HTML\n"
        "3. Загрузите обложку (фото)\n"
        "4. Добавьте YouTube/Vimeo ссылку\n"
        "5. Заполните SEO поля\n"
        "6. Нажмите «Сохранить»\n"
        "7. Статус: Черновик / Опубликована",
        f"{SC}/sc_blog.png", RGBColor(0x7c,0x3a,0xed))

    content_slide(prs, "SEO для статей", [
        "## Поля SEO",
        "SEO Title — заголовок для Google (до 60 символов)",
        "SEO Description — описание (до 160 символов)",
        "Keywords — ключевые слова через запятую",
        "",
        "## Автоматически генерируется:",
        "Open Graph теги для соцсетей",
        "Canonical URL",
        "Запись в sitemap.xml",
        "",
        "Проверяйте через: Google Search Console",
    ], f"{SC}/sc_article.png", RGBColor(0x7c,0x3a,0xed))

    section_slide(prs, 5, "Уведомления и Railway", RGBColor(0x7c,0x3a,0xed))

    content_slide(prs, "Когда и кому приходят уведомления", [
        "## Клиент получает:",
        "✅ Новая заявка принята (со ссылкой на ЛК)",
        "✅ Назначена няня (с датами)",
        "",
        "## Няня получает:",
        "🆕 Новое назначение (данные клиента, даты)",
        "⏰ Напоминание за 2 часа до смены",
        "",
        "## Администратор получает:",
        "🆕 Каждая новая заявка с контактами",
        "✅ Подтверждение смен",
    ], f"{SC}/sc_tg_entry.png", RGBColor(0x7c,0x3a,0xed))

    content_slide(prs, "Railway — переменные окружения", [
        "## Обязательные переменные:",
        "TELEGRAM_BOT_TOKEN — токен бота от @BotFather",
        "ADMIN_IDS — Telegram ID администраторов через запятую",
        "FLASK_SECRET_KEY — случайная строка 32+ символа",
        "ADMIN_TOKEN — пароль для dev-доступа",
        "",
        "## Опциональные:",
        "FORCE_HTTPS=1 — принудительный HTTPS",
        "STORAGE=sql — режим базы данных (по умолчанию JSON)",
    ], None, RGBColor(0x7c,0x3a,0xed))

    tip_slide(prs, "Чек-лист перед запуском", [
        "TELEGRAM_BOT_TOKEN и ADMIN_IDS выставлены в Railway",
        "Домен подключён и SSL работает",
        "BotFather: Menu Button URL = ваш-домен/app",
        "Тестовая заявка прошла и пришло уведомление",
        "Договоры и тарифы загружены в /static/contracts/",
        "Создана первая статья в блоге",
    ], RGBColor(0x7c,0x3a,0xed))

    prs.save(f"{OUT}/guide_admin.pptx")
    print("✅ guide_admin.pptx")


make_client()
make_nanny()
make_admin()
print("\nAll presentations done!")
